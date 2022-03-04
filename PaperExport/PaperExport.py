# !/usr/bin/env python
# -*- coding: utf-8 -*-
# @Author: gyy

from docx import Document
from docx.document import Document as _Document
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
import copy
import os


class PaperExport:
    """ Word、PPT导出修改程序 """
    def __init__(self):
        self.TYPE = {
            "1": ".docx",
            "2": ".ppt",
            "3": ".pdf",
            "4": ".xls"
        }
        self.type = ""
        self.file = ""

    def find_file(self, type):
        for root, dirs, files in os.walk("."):
            for f in files:
                if type in f:
                    return f

    def set_font(self, run, font_name, ppt=False, size=10.5):
        run.font.size = Pt(size)
        if ppt:
            run.font.name = font_name
            run.font.name = u"Times New Roman"
        else:
            run.font.name = u"Times New Roman"
            run.element.rPr.rFonts.set(qn('w:eastAsia'), font_name)

    def iter_block_items(self, parent):
        if isinstance(parent, _Document):
            parent_elm = parent.element.body
            # print(parent_elm.xml)
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            raise ValueError("something's not right")
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    def get_cell_text(self, cell, listen=False):
        cell_text = ""
        for block in self.iter_block_items(cell):
            if isinstance(block, Paragraph):
                cell_text += block.text.replace("卐卐", " " * 2).strip() + "\n"
            elif isinstance(block, Table):
                for row in block.rows:
                    for ce in row.cells:
                        if listen:
                            cell_text += ce.text + "\n"
                        else:
                            cell_text += "\n" + ce.text + "\n"
        if listen:
            cell_text = cell_text.replace("\n\n", "\n")
        return cell_text

    def read_word(self):
        res = []
        big_font = False
        listen = False
        reading = False
        doc = Document(self.find_file(".docx"))
        for t in range(len(doc.tables)):
            for r in range(len(doc.tables[t].rows)):
                for cell in doc.tables[t].rows[r].cells:
                    if cell.text or cell.tables:
                        rgb = ""
                        try:
                            rgb = str(cell.paragraphs[0].runs[0].font.color.rgb)
                        except:
                            pass
                        if rgb == "08AFA8":
                            if "选择题" in cell.text or "听力" in cell.text:
                                big_font = True
                                if "听力" in cell.text:
                                    listen = True
                            if "阅读" in cell.text:
                                reading = True
                        elif listen:
                            res.append(self.get_cell_text(cell))
                        elif "答案" not in cell.text and "未分类习题" not in cell.text and t != 0:
                            cell_text = self.get_cell_text(cell)
                            res.append(cell_text)
        if reading:
            big_font = False
            listen = False
        return res, big_font, listen

    def export_word(self):
        from docx.shared import Cm, Pt
        abs_path = os.path.abspath(self.file)
        doc = Document(abs_path)
        try:
            if "Print" in doc.paragraphs[0].runs[0].text:
                doc.paragraphs[0].runs[0].clear()
        except:
            pass
        # 边距、装订线
        for section in doc.sections:
            section.top_margin = Cm(2.54)
            section.bottom_margin = Cm(2.54)
            section.left_margin = Cm(3.17)
            section.right_margin = Cm(3.17)
            section.gutter = Cm(0)
        # 题型间添加行
        for t in range(len(doc.tables)):
            doc.tables[t].autofit = True
            doc.tables[t].columns[0].width = Cm(14.66)
            for r in range(len(doc.tables[t].rows)):
                for cell in doc.tables[t].rows[r].cells:
                    if cell.text or cell.tables:
                        # print(cell.text)
                        try:
                            rgb = cell.paragraphs[0].runs[0].font.color.rgb
                            if str(rgb) == "08AFA8":
                                for run in cell.paragraphs[0].runs:
                                    if run.text == "":
                                        run.clear()
                            # doc.tables[t - 1].add_row()
                        except:
                            pass
                        if t == 0 or "未分类习题" in cell.text:
                            continue
                        for tab in cell.tables:
                            for ro in tab.rows:
                                for ce in ro.cells:
                                    for para in ce.paragraphs:
                                        para.paragraph_format.line_spacing = 1.25
                                        for run in para.runs:
                                            self.set_font(run, u'汉仪书宋二简')

                        for para in cell.paragraphs:
                            para.paragraph_format.line_spacing = 1.25
                            para.paragraph_format.space_after = Pt(0)
                            for run in para.runs:
                                self.set_font(run, u'汉仪书宋二简')

                        if "答案" in cell.text or "解析" in cell.text:
                            for ce in doc.tables[t].rows[r].cells:
                                for para in ce.paragraphs:
                                    for run in para.runs:
                                        self.set_font(run, u'楷体')

        doc.save(abs_path)

    @staticmethod
    def cut_text(text, size=14):
        # 计算字符实际宽度
        def get_char_width(o):
            """Return the screen column width for unicode ordinal o."""
            widths = [
                (126, 1), (159, 0), (687, 1), (710, 0), (711, 1),
                (727, 0), (733, 1), (879, 0), (1154, 1), (1161, 0),
                (4347, 1), (4447, 2), (7467, 1), (7521, 0), (8369, 1),
                (8426, 0), (9000, 1), (9002, 2), (11021, 1), (12350, 2),
                (12351, 1), (12438, 2), (12442, 0), (19893, 2), (19967, 1),
                (55203, 2), (63743, 1), (64106, 2), (65039, 1), (65059, 0),
                (65131, 2), (65279, 1), (65376, 2), (65500, 1), (65510, 2),
                (120831, 1), (262141, 2), (1114109, 1),
            ]
            if o == 0xe or o == 0xf:
                return 0
            for num, wid in widths:
                if o <= num:
                    return wid
            return 1

        def get_line_width(line):
            wid = 0
            for char in line:
                wid += get_char_width(ord(char))
            return wid

        def cut_chn_line(line, wid):
            res = []
            if len(line) < wid / 2:
                res.append(line)
                return res
            line_len = 0
            temp = ""
            count = len(line)
            for char in line:
                temp += char
                line_len += get_char_width(ord(char))
                if line_len >= wid:
                    res.append(temp)
                    temp = ""
                    line_len = 0
                count -= 1
                if count == 0:
                    res.append(temp)
            return res

        text_list = text.splitlines()
        while '' in text_list:
            text_list.remove('')
        width = 22 * 72 // 16
        small_width = 14 * 72 // 16

        from textwrap import wrap
        def wrap_text(text_list, width):
            wrap_list = []
            for i in text_list:
                wrap_text = wrap(i, width)
                for t in wrap_text:
                    line_w = get_line_width(t)
                    if line_w > width:
                        cut = cut_chn_line(t, width - 20)
                        wrap_list.extend(cut)
                    else:
                        wrap_list.append(t)
            return wrap_list

        wrap_list = wrap_text(text_list, width)

        listen = False
        big_font = False
        if len(wrap_list) <= size * 0.5:
            listen = True
            big_font = True
            wrap_list = wrap_text(text_list, small_width)

        res_list = []
        paras = len(wrap_list) // size
        if paras > 0:
            for i in range(paras):
                res_list.append(wrap_list[i*size: (i+1)*size])
            if len(wrap_list) % size != 0:
                res_list.append(wrap_list[paras*size:])
        else:
            res_list.append(wrap_list)

        return res_list, listen, big_font

    def change_ppt_title(self, prs, slide):
        number = prs.slides[slide].shapes[0]
        title = prs.slides[slide].shapes[1]
        # 修改字体大小
        number.text_frame.paragraphs[0].runs[0].font.size = Pt(80)
        title.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
        # 修改文本框位置
        number.top = Cm(2)
        title.top = Cm(6)

    def change_ppt_content(self, prs, slide, text, listen=False, big_font=False):
        current_slide = slide
        text_list, listen, big_font = self.cut_text(text)
        for i in range(len(text_list)):
            # 复制幻灯片
            if i > 0:
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                temp = prs.slides[current_slide]
                new = prs.slides.add_slide(prs.slide_layouts[0])
                # print(new.background.fill)
                for shp in temp.shapes:
                    el = shp.element
                    newel = copy.deepcopy(el)
                    new.shapes._spTree.insert_element_before(newel, 'p:extLst')
                pic = new.shapes.add_picture('backgroud_image_xqf.png', 0, 0,
                                               height=prs.slide_height)
                # try:
                #     for _, value in six.iteritems(temp.part.rels):
                #         # Make sure we don't copy a notesSlide relation as that won't exist
                #         if "notesSlide" not in value.reltype:
                #             new.part.rels.add_relationship(value.reltype,
                #                                                     value._target,
                #                                                     value.rId)
                # except Exception as e:
                #     print(e)
                current_slide += 1
                xml_slides.remove(slides[-1])
                xml_slides.insert(current_slide, slides[-1])
            # 清空内容
            try:
                shapes = prs.slides[current_slide].shapes
                shapes.element.remove(shapes[0].element)
            except:
                pass
            # 添加文本框
            txt_box = prs.slides[current_slide].shapes.add_textbox(Cm(1), Cm(1.2), Cm(23.5), Cm(13))
            tf = txt_box.text_frame
            p = tf.add_paragraph()
            p.text = "\n".join(text_list[i])
            # 修改格式
            if listen:
                for run in p.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    self.set_font(run, u'汉仪书宋二简', ppt=True, size=24)
            else:
                p.font.color.rgb = RGBColor(255, 255, 255)
                if big_font:
                    self.set_font(p, u'汉仪书宋二简', ppt=True, size=24)
                else:
                    self.set_font(p, u'汉仪书宋二简', ppt=True, size=16)

        return current_slide - slide + 1

    def export_ppt(self):
        prs = Presentation(self.file)
        title_n = len(prs.slides)
        print(title_n)
        text_list, big_font, listen = self.read_word()
        # for i in text_list:
        #     print("{}: {}".format(text_list.index(i), i))
        content_slide = 0
        next_slide = 0
        flag = 0
        while content_slide < title_n / 2:
            # 修改标题页
            if flag == 0:
                self.change_ppt_title(prs, next_slide)
                next_slide += 1
                flag = 1
            # 修改内容页
            elif flag == 1:
                # 不加strip可能会无法修改格式
                text = text_list[content_slide].strip()
                add_slide = self.change_ppt_content(prs, next_slide, text, listen, big_font)
                next_slide += add_slide
                content_slide += 1
                flag = 0

        prs.save(self.file)

    def find_func(self):
        func_map = {
            self.TYPE["1"]: self.export_word,
            self.TYPE["2"]: self.export_ppt,
        }
        return func_map[self.type]()

    def run(self, type):
        self.type = self.TYPE[type]
        self.file = self.find_file(self.type)
        self.find_func()


def main():
    pe = PaperExport()
    pe.run("1")
    pe.run("2")
    input("运行完成")


main()
